"""
Document Text Extraction Service

Extracts text content from various document formats (PDF, DOCX, PPTX, TXT)
with page/slide number tracking for source attribution.
"""

from typing import Dict, List, Optional
from pathlib import Path
import structlog

logger = structlog.get_logger()


class ExtractionResult:
    """
    Container for extracted text with source attribution metadata.

    Attributes:
        text: Full extracted text
        pages: List of (page_number, text) tuples for page-level attribution
        metadata: Additional metadata (page count, author, etc.)
    """

    def __init__(
        self,
        text: str,
        pages: List[tuple[int, str]],
        metadata: Optional[Dict[str, str]] = None,
    ):
        self.text = text
        self.pages = pages
        self.metadata = metadata or {}

    def __repr__(self) -> str:
        return f"ExtractionResult(pages={len(self.pages)}, chars={len(self.text)})"


class DocumentExtractor:
    """
    Document text extraction service.

    Supports multiple formats and provides page-level granularity for
    source attribution requirements.
    """

    def __init__(self):
        """Initialize document extractor."""
        self.supported_extensions = {".pdf", ".docx", ".pptx", ".txt"}

    def is_supported(self, file_path: str) -> bool:
        """
        Check if file format is supported.

        Args:
            file_path: Path to document file

        Returns:
            True if format is supported
        """
        return Path(file_path).suffix.lower() in self.supported_extensions

    def extract(self, file_path: str) -> ExtractionResult:
        """
        Extract text from document based on file extension.

        Args:
            file_path: Path to document file

        Returns:
            ExtractionResult with text and page-level metadata

        Raises:
            ValueError: If file format is not supported
            Exception: If extraction fails
        """
        path = Path(file_path)
        extension = path.suffix.lower()

        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        if extension not in self.supported_extensions:
            raise ValueError(
                f"Unsupported file format: {extension}. "
                f"Supported formats: {', '.join(self.supported_extensions)}"
            )

        logger.info(
            "extracting_document",
            file_path=file_path,
            extension=extension,
        )

        try:
            if extension == ".pdf":
                return self._extract_pdf(file_path)
            elif extension == ".docx":
                return self._extract_docx(file_path)
            elif extension == ".pptx":
                return self._extract_pptx(file_path)
            elif extension == ".txt":
                return self._extract_txt(file_path)
            else:
                raise ValueError(f"Unsupported extension: {extension}")

        except Exception as e:
            logger.error(
                "extraction_failed",
                file_path=file_path,
                extension=extension,
                error=str(e),
                error_type=type(e).__name__,
            )
            raise

    def _extract_pdf(self, file_path: str) -> ExtractionResult:
        """
        Extract text from PDF file using pdfplumber.

        Args:
            file_path: Path to PDF file

        Returns:
            ExtractionResult with page-level text
        """
        try:
            import pdfplumber
        except ImportError:
            raise ImportError(
                "pdfplumber is required for PDF extraction. "
                "Install with: pip install pdfplumber"
            )

        pages = []
        all_text = []

        with pdfplumber.open(file_path) as pdf:
            metadata = {
                "page_count": str(len(pdf.pages)),
                "format": "PDF",
            }

            # Extract PDF metadata if available
            if pdf.metadata:
                if pdf.metadata.get("Author"):
                    metadata["author"] = pdf.metadata.get("Author")
                if pdf.metadata.get("Title"):
                    metadata["title"] = pdf.metadata.get("Title")

            for page_num, page in enumerate(pdf.pages, start=1):
                try:
                    text = page.extract_text() or ""
                    if text.strip():
                        pages.append((page_num, text))
                        all_text.append(text)
                except Exception as e:
                    logger.warning(
                        "pdf_page_extraction_failed",
                        file_path=file_path,
                        page_num=page_num,
                        error=str(e),
                    )
                    # Continue with other pages
                    pages.append((page_num, ""))

        full_text = "\n\n".join(all_text)

        logger.info(
            "pdf_extracted",
            file_path=file_path,
            pages=len(pages),
            chars=len(full_text),
        )

        return ExtractionResult(text=full_text, pages=pages, metadata=metadata)

    def _extract_docx(self, file_path: str) -> ExtractionResult:
        """
        Extract text from DOCX file using python-docx.

        Args:
            file_path: Path to DOCX file

        Returns:
            ExtractionResult with paragraph-level text
        """
        try:
            from docx import Document
        except ImportError:
            raise ImportError(
                "python-docx is required for DOCX extraction. "
                "Install with: pip install python-docx"
            )

        doc = Document(file_path)
        pages = []
        all_text = []

        # DOCX doesn't have explicit pages, so we treat sections as "pages"
        # For better granularity, we could count paragraphs or characters
        current_page = 1
        page_text = []
        paragraphs_per_page = 10  # Arbitrary grouping for source attribution

        for i, paragraph in enumerate(doc.paragraphs, start=1):
            text = paragraph.text.strip()
            if text:
                page_text.append(text)
                all_text.append(text)

            # Group paragraphs into "pages" for source attribution
            if i % paragraphs_per_page == 0 and page_text:
                pages.append((current_page, "\n".join(page_text)))
                page_text = []
                current_page += 1

        # Add remaining paragraphs
        if page_text:
            pages.append((current_page, "\n".join(page_text)))

        full_text = "\n\n".join(all_text)

        metadata = {
            "page_count": str(len(pages)),
            "paragraph_count": str(len(doc.paragraphs)),
            "format": "DOCX",
        }

        # Extract core properties if available
        if doc.core_properties:
            if doc.core_properties.author:
                metadata["author"] = doc.core_properties.author
            if doc.core_properties.title:
                metadata["title"] = doc.core_properties.title

        logger.info(
            "docx_extracted",
            file_path=file_path,
            pages=len(pages),
            paragraphs=len(doc.paragraphs),
            chars=len(full_text),
        )

        return ExtractionResult(text=full_text, pages=pages, metadata=metadata)

    def _extract_pptx(self, file_path: str) -> ExtractionResult:
        """
        Extract text from PPTX file using python-pptx.

        Args:
            file_path: Path to PPTX file

        Returns:
            ExtractionResult with slide-level text
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx is required for PPTX extraction. "
                "Install with: pip install python-pptx"
            )

        prs = Presentation(file_path)
        pages = []
        all_text = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []

            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

                # Also check for text in tables
                if hasattr(shape, "table"):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                slide_text.append(cell.text.strip())

            text = "\n".join(slide_text)
            if text:
                pages.append((slide_num, text))
                all_text.append(text)
            else:
                # Include empty slides for proper numbering
                pages.append((slide_num, ""))

        full_text = "\n\n".join(all_text)

        metadata = {
            "slide_count": str(len(prs.slides)),
            "format": "PPTX",
        }

        # Extract core properties if available
        if prs.core_properties:
            if prs.core_properties.author:
                metadata["author"] = prs.core_properties.author
            if prs.core_properties.title:
                metadata["title"] = prs.core_properties.title

        logger.info(
            "pptx_extracted",
            file_path=file_path,
            slides=len(pages),
            chars=len(full_text),
        )

        return ExtractionResult(text=full_text, pages=pages, metadata=metadata)

    def _extract_txt(self, file_path: str) -> ExtractionResult:
        """
        Extract text from TXT file.

        Args:
            file_path: Path to TXT file

        Returns:
            ExtractionResult with line-based text
        """
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()

        # Split into "pages" based on line count for source attribution
        lines = content.split("\n")
        lines_per_page = 50  # Arbitrary grouping
        pages = []

        for i in range(0, len(lines), lines_per_page):
            page_num = (i // lines_per_page) + 1
            page_text = "\n".join(lines[i : i + lines_per_page])
            if page_text.strip():
                pages.append((page_num, page_text))

        metadata = {
            "line_count": str(len(lines)),
            "page_count": str(len(pages)),
            "format": "TXT",
        }

        logger.info(
            "txt_extracted",
            file_path=file_path,
            lines=len(lines),
            pages=len(pages),
            chars=len(content),
        )

        return ExtractionResult(text=content, pages=pages, metadata=metadata)


# Factory function for dependency injection
def get_document_extractor() -> DocumentExtractor:
    """
    Get document extractor instance.

    Returns:
        DocumentExtractor instance
    """
    return DocumentExtractor()
